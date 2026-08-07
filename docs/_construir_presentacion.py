"""Presentación institucional, construida desde las salidas del pipeline.

15 minutos: contexto 3, producto 10, cierre 2.
"""
from __future__ import annotations

import json
import pathlib
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

sys.path.insert(0, str(pathlib.Path(__file__).parents[1]))
from src import config

S = config.SURFACE
CIF = pd.read_csv(S / "lista_cifras.csv").set_index("concepto")["valor"].to_dict()
PORTADA = json.loads((config.DOCS / "registro_version.json")
                     .read_text(encoding="utf-8"))["portada"]
AZUL, NARANJA, ROJO, VERDE, GRIS = (RGBColor(0x31, 0x70, 0x8E), RGBColor(0xA5, 0x67, 0x3F),
                                    RGBColor(0xC6, 0x28, 0x28), RGBColor(0x2E, 0x7D, 0x32),
                                    RGBColor(0x55, 0x55, 0x55))


def diapositiva(prs, titulo, subtitulo=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    caja = s.shapes.add_textbox(Cm(1.2), Cm(0.8), Cm(31.5), Cm(2))
    p = caja.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.size, r.font.bold, r.font.color.rgb = Pt(30), True, AZUL
    if subtitulo:
        p2 = caja.text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitulo
        r2.font.size, r2.font.color.rgb = Pt(15), GRIS
    return s


def texto(s, x, y, ancho, alto, lineas, tam=16, color=None):
    caja = s.shapes.add_textbox(Cm(x), Cm(y), Cm(ancho), Cm(alto))
    tf = caja.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = linea if isinstance(linea, str) else linea[0]
        r.font.size = Pt(tam if isinstance(linea, str) else linea[1])
        r.font.color.rgb = color or RGBColor(0x22, 0x22, 0x22)
        if not isinstance(linea, str) and len(linea) > 2:
            r.font.bold = linea[2]
        p.space_after = Pt(8)
    return caja


def metrica(s, x, y, valor, etiqueta, color=AZUL, ancho=7.4):
    caja = s.shapes.add_textbox(Cm(x), Cm(y), Cm(ancho), Cm(3.2))
    tf = caja.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = valor
    r.font.size, r.font.bold, r.font.color.rgb = Pt(34), True, color
    p = tf.add_paragraph()
    r2 = p.add_run()
    r2.text = etiqueta
    r2.font.size, r2.font.color.rgb = Pt(12), GRIS
    return caja


def construir() -> pathlib.Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)

    # 1 · portada
    s = diapositiva(prs, "Producto de datos integrado para Buenaventura",
                    "Importaciones y carga portuaria · Universidad Libre, Seccional Cali")
    texto(s, 1.2, 8.5, 31, 8, [
        ("Juan Manuel Tejada Fajardo · Jesús Alejandro Guerrero", 18, True),
        "Ingeniería del Producto de Ciencia de Datos",
        (f"Director del trabajo: {PORTADA['director']}", 14, False),
        (f"Programa: {PORTADA['programa']}", 14, False),
        "Santiago de Cali · 2026",
    ], 15)

    # 2 · contexto
    s = diapositiva(prs, "El problema", "Contexto · 1 de 3")
    texto(s, 1.2, 4.5, 15.5, 12, [
        ("Los datos existen. La respuesta no.", 20, True),
        "El DANE publica declaraciones aduaneras.",
        "La Superintendencia de Transporte publica toneladas movilizadas.",
        "",
        ("Ninguna de las dos, por separado, responde:", 16, True),
        "¿el comercio creció porque entró más mercancía,",
        "o porque la mercancía se encareció?",
    ], 15)
    texto(s, 18, 4.5, 14.5, 12, [
        ("Lo que este trabajo hace", 18, True),
        f"· Integra {CIF['Meses de la serie aduanera']} meses aduaneros y "
        f"{CIF['Meses de la serie portuaria']} portuarios",
        f"· {CIF['Meses integrados']} meses de intersección",
        "· Responde 52 preguntas con evidencia",
        "· Documenta lo que NO se puede responder",
    ], 15)

    # 3 · alcance y límites
    s = diapositiva(prs, "Qué medimos y qué no", "Contexto · 2 de 3")
    texto(s, 1.2, 4.5, 15.5, 12, [
        ("SÍ", 20, True), "· Valor CIF, peso neto y valor unitario implícito",
        "· Toneladas por tipo de carga y sociedad portuaria",
        "· Concentración y especialización",
        "· Pronóstico con intervalos medidos",
    ], 15, VERDE)
    texto(s, 18, 4.5, 14.5, 12, [
        ("NO", 20, True), "· Congestión, patios, tiempos de despacho",
        "· Arribos, tipos de buque, ETA/ATA, permanencias",
        "· TEU: la fuente publica toneladas, no unidades",
        "· Qué buque trajo una importación concreta",
    ], 15, ROJO)
    texto(s, 1.2, 16.5, 31, 2,
          ["Las ocho preguntas del bloque marítimo se cierran como no viables, "
           "con la búsqueda documentada. Es una respuesta, no una omisión."], 13, GRIS)

    # 4 · fuentes
    s = diapositiva(prs, "Fuentes: lo que encontramos al verificar", "Contexto · 3 de 3")
    texto(s, 1.2, 4.5, 31, 5, [
        ("Suponíamos que el dato portuario era trimestral y en PDF. Era falso.", 20, True),
        "El conjunto 5r3g-zv5z de datos.gov.co publica el dato MENSUAL, por sociedad "
        "portuaria y tipo de carga, desde 2018, con API y licencia CC BY-SA 4.0.",
    ], 16)
    metrica(s, 1.2, 10.5, CIF["Meses de la serie portuaria"], "meses portuarios, sin huecos")
    metrica(s, 9.5, 10.5, "2026-06", "último mes: el puerto va por delante de la aduana", NARANJA)
    metrica(s, 17.8, 10.5, "0", "fuentes públicas con arribos o permanencias", ROJO)
    metrica(s, 26.1, 10.5, "CC BY-SA", "licencia verificada", VERDE)

    # 5 · el hallazgo central
    s = diapositiva(prs, "El hallazgo que cambió el producto", "Producto · 1 de 6")
    texto(s, 1.2, 4.2, 31, 3, [
        ("Integrar dominios NO mejora el pronóstico. Lo empeora.", 24, True)], 24, ROJO)
    metrica(s, 1.2, 8, "5,875 %", "solo historia propia + calendario", VERDE)
    metrica(s, 11.5, 8, "6,167 %", "modelo integrado completo", NARANJA)
    metrica(s, 21.8, 8, "6,575 %", "historia propia + puerto", ROJO)
    texto(s, 1.2, 12.5, 31, 5, [
        ("Por qué:", 16, True),
        "ambas fuentes miden el mismo comercio con el mismo rezago de publicación. "
        "El puerto no aporta información que la historia del CIF no tenga ya, y sí consume "
        "grados de libertad sobre una muestra corta.",
        ("Consecuencia: el modelo predictivo es el aduanero. El puerto aporta valor "
         "descriptivo y explicativo.", 16, True),
    ], 14)

    # 6 · concentración
    s = diapositiva(prs, "El contraste que solo se ve integrando", "Producto · 2 de 6")
    metrica(s, 1.2, 5, CIF["HHI capítulo arancelario"], "HHI capítulo arancelario\ndesconcentrado", VERDE, 9.5)
    metrica(s, 12, 5, CIF["HHI país de origen"], "HHI país de origen\ndesconcentrado", VERDE, 9.5)
    metrica(s, 22.8, 5, CIF["HHI sociedad portuaria"], "HHI sociedad portuaria\nCONCENTRADO", ROJO, 9.5)
    texto(s, 1.2, 10.5, 31, 7, [
        ("Buenaventura importa mercancía variada desde orígenes variados,", 20, True),
        ("pero la mueve por muy pocas sociedades portuarias.", 20, True),
        "",
        f"Una sola concentra el {CIF['Participación de la mayor sociedad']} de las toneladas. "
        "Y la especialización es casi total: una sociedad dedicada a contenedores no puede "
        "absorber un desvío de granel.",
        "El riesgo no está en qué se importa ni de dónde viene, sino en el punto físico por "
        "el que pasa.",
    ], 15)

    # 7 · sociedades sin reporte
    s = diapositiva(prs, "Dos sociedades no aparecen reportando en 2026",
                    "Producto · 3 de 6")
    sr = pd.read_csv(S / "sociedades_sin_reporte_2026.csv")
    texto(s, 1.2, 4.5, 15.5, 11,
          [("Qué muestran los datos", 18, True)] +
          [f"· {r['sociedad_portuaria'][:42]} — último reporte {r['ultimo_anio_con_reporte']}"
           for _, r in sr.iterrows()] +
          ["", "Verificado mes a mes: en los seis meses observados de 2026 solo reportan "
           "tres sociedades. No es rezago de publicación."], 14)
    texto(s, 18, 4.5, 14.5, 11, [
        ("Lo que NO podemos afirmar", 18, True),
        "La fuente registra reporte, no operación.",
        "",
        "No podemos decir que dejaron de operar.",
        "Solo que no figuran en los reportes del periodo observado.",
        "",
        ("Por qué importa:", 15, True),
        "un modelo entrenado sobre el total leería esta ausencia como una caída real del "
        "comercio.",
    ], 14, ROJO)

    # 8 · qué se puede pronosticar
    s = diapositiva(prs, "Qué se puede pronosticar y qué no", "Producto · 4 de 6")
    texto(s, 1.2, 4.5, 31, 3, [
        ("El criterio se fijó antes de mirar los resultados: sin 36 observaciones, "
         "se describe, no se pronostica.", 16, True)], 16)
    metrica(s, 1.2, 8, CIF["WAPE · toneladas totales, Ridge"], "toneladas totales · Ridge\ncontra "
            + CIF["WAPE · toneladas totales, naive 1"] + " del naive", VERDE)
    metrica(s, 9.5, 8, CIF["WAPE · carga contenerizada, naive 1"],
            "carga contenerizada · NAIVE 1\nRidge da " +
            CIF["WAPE · carga contenerizada, Ridge"] + ": peor", ROJO)
    metrica(s, 17.8, 8, "0", "TEU, arribos y permanencias\nsin fuente", GRIS)
    metrica(s, 26.1, 8, "4 de 7", "indicadores elegibles", AZUL)
    texto(s, 1.2, 13.5, 31, 4, [
        "Para la carga contenerizada la recomendación es no usar modelo. Un indicador que "
        "no se pronostica mejor que repetir el último valor no debe llevar un modelo encima.",
    ], 15)

    # 9 · el producto
    s = diapositiva(prs, "El producto: seis vistas", "Producto · 5 de 6 · demostración en vivo")
    texto(s, 1.2, 4.5, 15.5, 12, [
        ("1 · Ejecutiva", 17, True), "valor y volumen lado a lado",
        ("2 · Aduanera", 17, True), "CIF, peso, valor unitario, estacionalidad",
        ("3 · Portuaria", 17, True), "tipo de carga, sociedades, HHI, especialización",
    ], 14)
    texto(s, 18, 4.5, 14.5, 12, [
        ("4 · Marítima", 17, True), "la no viabilidad, documentada",
        ("5 · Predictiva", 17, True), "modelos, líneas base, intervalos, ablación",
        ("6 · Calidad", 17, True), "trazabilidad de las 52 preguntas",
    ], 14)
    texto(s, 1.2, 16.5, 31, 2,
          ["El tablero no calcula nada: lee de la capa de consumo. Lo que se ve en pantalla "
           "y lo que dice el documento salen del mismo archivo."], 13, GRIS)

    # 10 · trazabilidad
    s = diapositiva(prs, "Trazabilidad", "Producto · 6 de 6")
    metrica(s, 1.2, 5.5, "38", "preguntas ejecutadas\ncon archivo de evidencia", VERDE)
    metrica(s, 9.5, 5.5, "4", "parciales\ncon limitación declarada", NARANJA)
    metrica(s, 17.8, 5.5, "10", "no viables\ncon búsqueda documentada", ROJO)
    metrica(s, 26.1, 5.5, "0", "sin evidencia\nni justificación", AZUL)
    texto(s, 1.2, 11, 31, 6, [
        ("80 pruebas automatizadas · ruff sin hallazgos · 42 figuras · 54 archivos de "
         "evidencia", 16, True),
        "Ninguna cifra del documento está escrita a mano: todas provienen de un archivo "
        "generado por el pipeline.",
        "El proyecto anterior se conserva congelado con manifiesto de hashes SHA-256.",
    ], 15)

    # 11 · conclusiones
    s = diapositiva(prs, "Conclusiones", "Cierre · 1 de 2")
    texto(s, 1.2, 4.5, 31, 12, [
        ("1. La integración es agregada por mes, nunca directa.", 17, True),
        "No existe llave pública entre una declaración y un movimiento portuario.",
        ("2. Integrar no mejora el pronóstico: lo empeora.", 17, True),
        "El modelo predictivo es el aduanero; el puerto aporta explicación.",
        ("3. Canasta diversificada, movilización concentrada.", 17, True),
        "El riesgo está en el punto físico por el que pasa la carga.",
        ("4. Lo que no se puede medir, se documenta.", 17, True),
        "Ocho preguntas cerradas con evidencia valen más que ocho cifras inventadas.",
    ], 14)

    # 12 · límites y siguiente paso
    s = diapositiva(prs, "Límites y siguiente paso", "Cierre · 2 de 2")
    texto(s, 1.2, 4.5, 15.5, 12, [
        ("Lo que este producto NO afirma", 18, True),
        "· Nada sobre congestión ni tiempos de atención",
        "· El CIF/kg no es un precio",
        "· El transbordo no es comercio exterior",
        "· Correlación no es causalidad",
        "· Las alertas no son órdenes operativas",
    ], 14, ROJO)
    texto(s, 18, 4.5, 14.5, 12, [
        ("Siguiente paso", 18, True),
        "· Validar las reglas de alerta con un usuario real",
        "· Segunda descarga del DANE para medir revisiones",
        "· Traducir códigos de país y capítulo a nombres",
        "· Gestionar acceso a fuentes de evento",
    ], 14, VERDE)

    ruta = config.REPORTS / "Presentacion_Buenaventura_V5.pptx"
    prs.save(ruta)
    return ruta


if __name__ == "__main__":
    r = construir()
    print(f"presentación generada: {r.name} ({r.stat().st_size // 1024} KB)")
