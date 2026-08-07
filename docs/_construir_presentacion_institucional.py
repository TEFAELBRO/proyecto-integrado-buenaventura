"""Presentación institucional de Universidad Libre para la sustentación V5.

Reutiliza los fondos, logos, colores y tipografía extraídos de la plantilla oficial
(reports/plantilla_institucional/). Todas las cifras se leen de lista_cifras.csv:
si un concepto no está verificado, la construcción falla en lugar de inventarlo.

Reparto de tiempo: contexto 3 min · producto 10 min · cierre 2 min.
"""
from __future__ import annotations

import json
import pathlib
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt

sys.path.insert(0, str(pathlib.Path(__file__).parents[1]))
from src import config  # noqa: E402

S = config.SURFACE
PLANTILLA = config.REPORTS / "plantilla_institucional"
FIGURAS = config.REPORTS / "figures"

_CIFRAS = pd.read_csv(S / "lista_cifras.csv").set_index("concepto")["valor"].astype(str).to_dict()
_USADAS: list[str] = []
REGISTRO = json.loads((config.DOCS / "registro_version.json").read_text(encoding="utf-8"))
PORTADA = REGISTRO["portada"]


def c(concepto: str) -> str:
    """Devuelve una cifra verificada. Falla si no está en el catálogo."""
    if concepto not in _CIFRAS:
        raise KeyError(f"cifra no verificada en lista_cifras.csv: {concepto!r}")
    _USADAS.append(concepto)
    return _CIFRAS[concepto]


def mil(concepto: str) -> str:
    """La misma cifra verificada, con separador de miles. No altera el valor."""
    entero, _, decimales = c(concepto).partition(",")
    signo, digitos = ("-", entero[1:]) if entero.startswith("-") else ("", entero)
    partes = []
    while len(digitos) > 3:
        partes.insert(0, digitos[-3:])
        digitos = digitos[:-3]
    partes.insert(0, digitos)
    return signo + ".".join(partes) + (f",{decimales}" if decimales else "")


# Paleta muestreada de la plantilla oficial
MARRON = RGBColor(0x84, 0x3C, 0x0C)
NARANJA = RGBColor(0xE0, 0x8C, 0x28)
DORADO = RGBColor(0xEE, 0xC1, 0x2E)
DURAZNO = RGBColor(0xF8, 0xD3, 0xBA)
ROSA = RGBColor(0xD3, 0xA2, 0x95)
CREMA = RGBColor(0xFA, 0xEC, 0xC0)
GRIS = RGBColor(0x76, 0x71, 0x71)
GRIS_CLARO = RGBColor(0xD9, 0xD9, 0xD9)
TINTA = RGBColor(0x3B, 0x38, 0x38)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x37, 0x6E, 0x3F)
ROJO = RGBColor(0xA6, 0x2B, 0x2B)

TITULAR = "Montserrat"
CUERPO = "Calibri"

ANCHO, ALTO = 33.867, 19.05
FECHA = "Agosto de 2026."


# --------------------------------------------------------------------------- base

def _fondo(prs, imagen: str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(str(PLANTILLA / imagen), 0, 0, Cm(ANCHO), Cm(ALTO))
    return s


def lamina(prs, titulo: str = "", subtitulo: str = "", bloque: str = ""):
    """Diapositiva de contenido con el fondo institucional y su encabezado."""
    s = _fondo(prs, "fondo_contenido.png")
    _texto(s, 1.55, 1.18, 7.4, 0.8, [(FECHA, 10, True)], color=BLANCO, fuente=TITULAR)
    if titulo:
        _texto(s, 1.7, 2.45, 24.0, 1.5, [(titulo, 25, True)], color=MARRON, fuente=TITULAR)
    if subtitulo:
        _texto(s, 1.75, 3.65, 27.0, 1.0, [(subtitulo, 12.5, False)], color=GRIS, fuente=CUERPO)
    if bloque:
        etq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Cm(27.6), Cm(2.55), Cm(4.7), Cm(0.95))
        _pinta(etq, DURAZNO, MARRON)
        _rellena_texto(etq, [(bloque, 9.5, True)], MARRON, CUERPO, PP_ALIGN.CENTER)
    s.shapes.add_picture(str(PLANTILLA / "logo_icono.png"), Cm(13.9), Cm(16.5),
                         Cm(2.4), Cm(2.35))
    return s


def _pinta(forma, relleno, borde=None, grosor=1.0):
    if relleno is None:
        forma.fill.background()
    else:
        forma.fill.solid()
        forma.fill.fore_color.rgb = relleno
    if borde is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borde
        forma.line.width = Pt(grosor)
    forma.shadow.inherit = False


def _lineas(tf, lineas, color, fuente, align=PP_ALIGN.LEFT, interlinea=1.0):
    tf.word_wrap = True
    for i, linea in enumerate(lineas):
        txt, tam, negrita, *resto = (linea if isinstance(linea, (list, tuple))
                                     else (linea, 14, False))
        col = resto[0] if resto else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = interlinea
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(tam)
        r.font.bold = negrita
        r.font.name = fuente
        r.font.color.rgb = col
        p.space_after = Pt(3)


def _rellena_texto(forma, lineas, color=TINTA, fuente=CUERPO, align=PP_ALIGN.CENTER,
                   margen=0.12):
    tf = forma.text_frame
    tf.margin_left = tf.margin_right = Cm(margen)
    tf.margin_top = tf.margin_bottom = Cm(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _lineas(tf, lineas, color, fuente, align)


def _texto(s, x, y, an, al, lineas, color=TINTA, fuente=CUERPO,
           align=PP_ALIGN.LEFT, interlinea=1.0):
    caja = s.shapes.add_textbox(Cm(x), Cm(y), Cm(an), Cm(al))
    caja.text_frame.margin_left = caja.text_frame.margin_top = 0
    _lineas(caja.text_frame, lineas, color, fuente, align, interlinea)
    return caja


def vinetas(s, x, y, an, items, tam=13, color=TINTA, punto="—"):
    lineas = []
    for it in items:
        txt, negrita = (it if isinstance(it, tuple) else (it, False))
        lineas.append((f"{punto}  {txt}", tam, negrita))
    return _texto(s, x, y, an, 1.0, lineas, color=color, interlinea=1.25)


# --------------------------------------------------- representaciones gráficas

def proceso(s, y, pasos, alto=2.3, x0=1.9, ancho_total=30.1, relleno=DURAZNO,
            borde=MARRON, tam=11.5):
    """SmartArt de proceso: cadena de galones encadenados."""
    n = len(pasos)
    an = ancho_total / n
    for i, paso in enumerate(pasos):
        forma = s.shapes.add_shape(MSO_SHAPE.CHEVRON if i else MSO_SHAPE.PENTAGON,
                                   Cm(x0 + i * an), Cm(y), Cm(an + 0.55), Cm(alto))
        _pinta(forma, relleno, borde)
        titulo, detalle = (paso if isinstance(paso, tuple) else (paso, ""))
        lineas = [(titulo, tam, True, MARRON)]
        if detalle:
            lineas.append((detalle, tam - 2.5, False, TINTA))
        _rellena_texto(forma, lineas, margen=0.5)


def ciclo(s, cx, cy, radio, pasos, r_nodo=2.15, tam=10.5, achatado=0.42):
    """SmartArt de ciclo: nodos dispuestos en círculo con el centro rotulado."""
    import math
    n = len(pasos)
    for i, paso in enumerate(pasos):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + radio * math.cos(ang) - r_nodo
        y = cy + radio * math.sin(ang) * achatado - r_nodo * 0.62
        nodo = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y),
                                  Cm(r_nodo * 2), Cm(r_nodo * 1.24))
        _pinta(nodo, DURAZNO if i % 2 == 0 else CREMA, MARRON)
        _rellena_texto(nodo, [(paso, tam, True, MARRON)], margen=0.18)
        ang2 = -math.pi / 2 + 2 * math.pi * (i + 0.5) / n
        fx = cx + (radio * 0.62) * math.cos(ang2)
        fy = cy + (radio * 0.62) * math.sin(ang2) * achatado
        punta = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(fx - 0.16), Cm(fy - 0.16),
                                   Cm(0.32), Cm(0.32))
        _pinta(punta, NARANJA, None)


def caja(s, x, y, an, al, lineas, relleno=DURAZNO, borde=MARRON,
         forma=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, margen=0.2):
    f = s.shapes.add_shape(forma, Cm(x), Cm(y), Cm(an), Cm(al))
    _pinta(f, relleno, borde)
    _rellena_texto(f, lineas, align=align, margen=margen)
    return f


def conector(s, x1, y1, x2, y2, color=MARRON, grosor=1.25):
    ln = s.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(grosor)
    return ln


def mapa_conceptual(s, centro, nodos, cx=16.9, cy=10.6, an_c=9.2, al_c=2.6,
                    radio_x=11.4, radio_y=3.9):
    """Nodo central con ramas rotuladas alrededor."""
    import math
    n = len(nodos)
    posiciones = []
    for i in range(n):
        ang = math.pi * (0.5 + 2 * i / n) if n > 2 else math.pi * (i + 0.5)
        posiciones.append((cx + radio_x * math.cos(ang), cy + radio_y * math.sin(ang)))
    for (px, py), nodo in zip(posiciones, nodos):
        titulo, detalle = (nodo if isinstance(nodo, tuple) else (nodo, ""))
        an, al = 8.4, 2.35
        conector(s, cx, cy, px, py)
        lineas = [(titulo, 11, True, MARRON)]
        if detalle:
            lineas.append((detalle, 9.5, False, TINTA))
        caja(s, px - an / 2, py - al / 2, an, al, lineas, CREMA, NARANJA)
    caja(s, cx - an_c / 2, cy - al_c / 2, an_c, al_c,
         [(centro, 14, True, BLANCO)], MARRON, MARRON)


def linea_tiempo(s, y, hitos, x0=2.2, ancho_total=29.5):
    """Línea de tiempo horizontal con hitos alternados."""
    barra = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x0), Cm(y - 0.22),
                               Cm(ancho_total), Cm(0.45))
    _pinta(barra, NARANJA, None)
    n = len(hitos)
    paso = ancho_total / n
    for i, (fecha, texto) in enumerate(hitos):
        px = x0 + paso * (i + 0.5)
        punto = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(px - 0.4), Cm(y - 0.62),
                                   Cm(0.8), Cm(0.8))
        _pinta(punto, MARRON, BLANCO, 1.5)
        arriba = i % 2 == 0
        cy = y - 3.35 if arriba else y + 0.95
        caja(s, px - paso / 2 + 0.25, cy, paso - 0.5, 2.4,
             [(fecha, 12, True, MARRON), (texto, 9.5, False, TINTA)], CREMA, NARANJA)
        conector(s, px, y - 0.62 if arriba else y + 0.62, px, cy + (2.4 if arriba else 0))


def arbol(s, raiz, ramas, y_raiz=5.6, y_ramas=9.4, x0=1.9, ancho_total=30.1):
    """Diagrama de árbol: una raíz que se abre en ramas."""
    an_r = 15.0
    cx = x0 + ancho_total / 2
    caja(s, cx - an_r / 2, y_raiz, an_r, 1.9, [(raiz, 13.5, True, BLANCO)], MARRON, MARRON)
    n = len(ramas)
    an = (ancho_total - 0.6 * (n - 1)) / n
    conector(s, cx, y_raiz + 1.9, cx, y_ramas - 1.0)
    for i, rama in enumerate(ramas):
        px = x0 + i * (an + 0.6)
        titulo, detalle = (rama if isinstance(rama, tuple) else (rama, ""))
        lineas = [(titulo, 11.5, True, MARRON)]
        if detalle:
            lineas.append((detalle, 9.5, False, TINTA))
        caja(s, px, y_ramas, an, 3.4, lineas, DURAZNO, NARANJA)
        conector(s, cx, y_ramas - 1.0, px + an / 2, y_ramas - 1.0)
        conector(s, px + an / 2, y_ramas - 1.0, px + an / 2, y_ramas)


def jerarquia(s, cima, hijos, y=5.4, x0=1.9, ancho_total=30.1):
    """SmartArt jerárquico: una cima y una fila de descendientes."""
    caja(s, x0 + ancho_total / 2 - 8.5, y, 17.0, 1.8,
         [(cima, 13.5, True, BLANCO)], MARRON, MARRON)
    n = len(hijos)
    an = (ancho_total - 0.55 * (n - 1)) / n
    y_h = y + 3.5
    cx = x0 + ancho_total / 2
    conector(s, cx, y + 1.8, cx, y_h - 0.9)
    for i, hijo in enumerate(hijos):
        px = x0 + i * (an + 0.55)
        num, titulo, detalle = hijo
        conector(s, cx, y_h - 0.9, px + an / 2, y_h - 0.9)
        conector(s, px + an / 2, y_h - 0.9, px + an / 2, y_h)
        f = caja(s, px, y_h, an, 4.6,
                 [(num, 15, True, NARANJA), (titulo, 11.5, True, MARRON),
                  (detalle, 9.5, False, TINTA)], BLANCO, MARRON)
        f.line.width = Pt(1.5)


def venn(s, izq, der, inter, cx=16.9, cy=11.0, r=5.4):
    a = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(cx - r - 2.5), Cm(cy - r * 0.72),
                           Cm(r * 2), Cm(r * 1.44))
    _pinta(a, DURAZNO, MARRON)
    a.fill.transparency = 0.35
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(cx - r + 2.5), Cm(cy - r * 0.72),
                           Cm(r * 2), Cm(r * 1.44))
    _pinta(b, CREMA, NARANJA)
    b.fill.transparency = 0.35
    _texto(s, cx - 8.4, cy - 2.3, 6.0, 4.6, izq, color=MARRON,
           align=PP_ALIGN.CENTER, interlinea=1.15)
    _texto(s, cx + 2.4, cy - 2.3, 6.0, 4.6, der, color=MARRON,
           align=PP_ALIGN.CENTER, interlinea=1.15)
    _texto(s, cx - 2.3, cy - 1.4, 4.6, 3.4, inter, color=TINTA,
           align=PP_ALIGN.CENTER, interlinea=1.15)


def mapa_red(s, centro, nodos, cx=16.9, cy=11.2, rx=11.6, ry=4.4):
    """Mapa de red: centro y nodos periféricos conectados entre sí."""
    import math
    n = len(nodos)
    pts = []
    for i in range(n):
        ang = 2 * math.pi * i / n - math.pi / 2
        pts.append((cx + rx * math.cos(ang), cy + ry * math.sin(ang)))
    for i, (px, py) in enumerate(pts):
        conector(s, cx, cy, px, py, GRIS_CLARO, 1.0)
        qx, qy = pts[(i + 1) % n]
        conector(s, px, py, qx, qy, GRIS_CLARO, 0.75)
    for (px, py), nodo in zip(pts, nodos):
        titulo, detalle = (nodo if isinstance(nodo, tuple) else (nodo, ""))
        lineas = [(titulo, 10.5, True, MARRON)]
        if detalle:
            lineas.append((detalle, 9, False, TINTA))
        caja(s, px - 3.8, py - 1.15, 7.6, 2.3, lineas, CREMA, NARANJA)
    caja(s, cx - 4.3, cy - 1.35, 8.6, 2.7, [(centro, 13, True, BLANCO)], MARRON, MARRON,
         forma=MSO_SHAPE.OVAL)


def rejilla(s, y, filas, alto=3.0, x0=1.9, ancho_total=30.1, sep=0.5,
            relleno=BLANCO, borde=MARRON):
    """Rejilla de tarjetas: filas de bloques del mismo alto."""
    for j, fila in enumerate(filas):
        n = len(fila)
        an = (ancho_total - sep * (n - 1)) / n
        for i, celda in enumerate(fila):
            titulo, *cuerpo = celda
            lineas = [(titulo, 11, True, MARRON)]
            lineas += [(t, 9.5, False, TINTA) for t in cuerpo]
            caja(s, x0 + i * (an + sep), y + j * (alto + sep), an, alto, lineas,
                 relleno, borde, align=PP_ALIGN.LEFT, margen=0.28)


def tabla(s, x, y, an, encabezados, filas, anchos=None, tam=10.5, alto_fila=0.72):
    n_f, n_c = len(filas) + 1, len(encabezados)
    t = s.shapes.add_table(n_f, n_c, Cm(x), Cm(y), Cm(an),
                           Cm(alto_fila * n_f)).table
    if anchos:
        total = sum(anchos)
        for j, w in enumerate(anchos):
            t.columns[j].width = Emu(int(Cm(an) * w / total))
    for j, h in enumerate(encabezados):
        celda = t.cell(0, j)
        celda.fill.solid()
        celda.fill.fore_color.rgb = MARRON
        _lineas(celda.text_frame, [(str(h), tam, True, BLANCO)], BLANCO, CUERPO)
        celda.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, fila in enumerate(filas, start=1):
        for j, v in enumerate(fila):
            celda = t.cell(i, j)
            celda.fill.solid()
            celda.fill.fore_color.rgb = BLANCO if i % 2 else CREMA
            negrita = isinstance(v, tuple)
            txt = v[0] if negrita else v
            col = v[1] if negrita and len(v) > 1 else TINTA
            _lineas(celda.text_frame, [(str(txt), tam, negrita, col)], TINTA, CUERPO)
            celda.vertical_anchor = MSO_ANCHOR.MIDDLE
    return t


def figura(s, nombre, x, y, an, al):
    ruta = FIGURAS / nombre
    if not ruta.exists():
        raise FileNotFoundError(f"falta la figura {ruta}")
    marco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x - 0.15), Cm(y - 0.15),
                               Cm(an + 0.3), Cm(al + 0.3))
    _pinta(marco, BLANCO, GRIS_CLARO)
    s.shapes.add_picture(str(ruta), Cm(x), Cm(y), Cm(an), Cm(al))


PIE = 15.82  # la franja inferior de la plantilla empieza aquí: nada la invade


def nota(s, y, texto, color=GRIS):
    """Nota al pie de la diapositiva, anclada por arriba de la franja institucional."""
    lineas = max(1, -(-len(texto) // 140))
    alto = lineas * 0.45 + 0.12
    return _texto(s, 1.9, min(y, PIE - alto), 30.1, alto, [(texto, 9.5, False)],
                  color=color, interlinea=1.15)


def notas_orador(s, texto):
    s.notes_slide.notes_text_frame.text = texto


# ------------------------------------------------------------------ construcción

def construir() -> pathlib.Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(ANCHO), Cm(ALTO)

    # 1 · Portada institucional -------------------------------------------------
    s = _fondo(prs, "fondo_portada.png")
    notas_orador(s, "Portada institucional. Saludo, nombre y programa. 10 segundos.")

    # 2 · Título, autores y director -------------------------------------------
    s = _fondo(prs, "fondo_contenido.png")
    _texto(s, 1.55, 1.18, 7.4, 0.8, [(FECHA, 10, True)], color=BLANCO, fuente=TITULAR)
    _texto(s, 1.85, 2.05, 20.0, 1.6, [("Ciencia de Datos", 28, True)],
           color=MARRON, fuente=TITULAR)
    _texto(s, 1.85, 3.55, 24.0, 1.0,
           [("“Tomando decisiones basados en conocimiento”", 17, True)],
           color=MARRON, fuente=TITULAR)
    campos = [
        ("Título", "Producto de datos para el análisis integrado del comercio de "
                   "importación y el tráfico portuario de Buenaventura", 6.6),
        ("Autores", " · ".join(PORTADA["integrantes"]), 9.8),
        ("Director", f"{PORTADA['director']}  ·  {PORTADA['institucion']}", 13.0),
    ]
    for etiqueta, valor, y in campos:
        conector(s, 7.7, y, 26.7, y, MARRON, 1.0)
        _texto(s, 7.75, y + 0.15, 19.0, 0.7, [(etiqueta, 11.5, False)], color=MARRON)
        _texto(s, 7.75, y + 0.85, 18.9, 1.6, [(valor, 13.5, True)], color=TINTA,
               interlinea=1.15)
    s.shapes.add_picture(str(PLANTILLA / "logo_icono.png"), Cm(13.9), Cm(16.5),
                         Cm(2.4), Cm(2.35))
    notas_orador(s, "Presentar el título y los integrantes. 15 segundos.")

    # 3 · Agenda ----------------------------------------------------------------
    s = lamina(prs, "Contenido de la sustentación",
               "Quince minutos repartidos en tres momentos", "15 minutos")
    rejilla(s, 5.5, [
        [("1 · Información general", "Introducción y procesos generales"),
         ("2 · Planteamiento del problema",
          "Descripción, antecedentes, formulación y sistematización"),
         ("3 · Objetivos", "General y específicos")],
        [("4 · Justificación, alcance y delimitación", "Por qué, hasta dónde y con qué"),
         ("5 · Marco de referencia", "Teórico, conceptual, contextual y legal"),
         ("6 · Metodología", "Fundamento de la ciencia de datos")],
    ], alto=2.6)
    rejilla(s, 11.5, [
        [("7 · Ingeniería y despliegue del producto",
          "Entorno, cuadro de mando y demostración en vivo"),
         ("8 · Modelo de negocio y comunicación", "Propuesta de valor y difusión"),
         ("9 · Cierre", "Conclusiones, recomendaciones y fuentes")],
    ], alto=2.4, relleno=DURAZNO)
    _texto(s, 1.9, 14.6, 30.1, 1.0,
           [("Contexto 3 min   ·   Producto de datos 10 min   ·   Cierre 2 min",
             12, True)], color=MARRON, align=PP_ALIGN.CENTER)
    notas_orador(s, "Anunciar el recorrido y el reparto de tiempo. 15 segundos.")

    # ===================== BLOQUE 1 · CONTEXTO (3 minutos) =====================

    # 4 · Introducción ----------------------------------------------------------
    s = lamina(prs, "Introducción",
               "El ciclo que recorre el proyecto, de la fuente oficial a la decisión",
               "Contexto")
    ciclo(s, 16.9, 10.4, 9.2, ["Fuentes\noficiales", "Ingesta y\ntrazabilidad",
                               "Análisis\nexploratorio", "Modelado y\nvalidación",
                               "Cuadro de\nmando", "Decisión\ninformada"])
    caja(s, 12.9, 9.15, 8.0, 2.5,
         [("Comercio de\nBuenaventura", 13, True, BLANCO)], MARRON, MARRON,
         forma=MSO_SHAPE.OVAL)
    nota(s, 15.75, f"Dos dominios oficiales integrados: {c('Meses de la serie aduanera')} "
                  f"meses de aduana ({c('Periodo aduanero')}) y "
                  f"{c('Meses de la serie portuaria')} meses de puerto "
                  f"({c('Periodo portuario')}).")
    notas_orador(s, "El proyecto no es un tablero: es un ciclo completo de datos. "
                    "Insistir en que las dos fuentes son oficiales y verificables. 25 s.")

    # 5 · Procesos generales ----------------------------------------------------
    s = lamina(prs, "Procesos generales",
               "Mapa de procesos: cuatro capas de datos con evidencia en cada paso",
               "Contexto")
    proceso(s, 5.6, [
        ("RAW", "Fuente inmutable\ncon hash SHA-256"),
        ("LANDING", "Normalización\nde formatos"),
        ("TRUSTED", "Series validadas\ny reconciliadas"),
        ("SURFACE", "Evidencia lista\npara consumo"),
        ("PRODUCTO", "Cuadro de mando\ny documento"),
    ], alto=3.0)
    rejilla(s, 9.6, [[
        ("Control de entrada", "Cada archivo se registra con su huella digital y su fecha "
                               "de descarga antes de ser tocado."),
        ("Control de proceso", "Reconciliación mes a mes contra el total publicado y "
                               "bitácora de defectos de la fuente."),
        ("Control de salida", f"{c('Pruebas automatizadas')} pruebas automatizadas y "
                              "revisión de estilo sin hallazgos."),
    ]], alto=3.3)
    nota(s, 13.6, "Ningún resultado del producto existe sin un archivo de evidencia que lo "
                  "respalde y sin la pregunta del catálogo que lo originó.")
    notas_orador(s, "Explicar que la arquitectura por capas es lo que hace auditable el "
                    "resultado. 25 segundos.")

    # 6 · Descripción del problema ---------------------------------------------
    s = lamina(prs, "Planteamiento del problema",
               "Mapa conceptual: qué impide hoy leer el comercio de Buenaventura como un "
               "solo fenómeno", "Contexto")
    mapa_conceptual(s, "La información existe,\npero está partida",
                    [("Dos registros separados",
                      "La aduana mide declaraciones; el puerto mide toneladas movilizadas."),
                     ("Sin llave común",
                      "No hay identificador público que una una declaración con un "
                      "movimiento de carga."),
                     ("Conceptos no equivalentes",
                      "El peso neto aduanero y las toneladas portuarias no miden lo mismo."),
                     ("Lectura incompleta",
                      "Con una sola fuente no se distingue si un mes cambió por valor o "
                      "por volumen.")],
                    cy=10.2)
    notas_orador(s, "El problema no es falta de datos: es que están partidos y nadie los "
                    "lee juntos. 25 segundos.")

    # 7 · Antecedentes ----------------------------------------------------------
    s = lamina(prs, "Antecedentes",
               "Línea de tiempo: de las fuentes disponibles al producto integrado",
               "Contexto")
    linea_tiempo(s, 10.6, [
        ("2012", "El DANE publica microdatos mensuales de importación por aduana."),
        ("2018", "Supertransporte abre el tráfico portuario mensual por sociedad."),
        ("V4", "Producto aduanero validado con backtesting y 80 pruebas."),
        ("2026", "V5 integra aduana y puerto y mide si la integración aporta."),
    ])
    nota(s, 15.0, "El hallazgo que abrió la versión 5: el dataset portuario resultó ser "
                  "mensual y descargable por interfaz de programación, no trimestral en "
                  "PDF como suponía el diagnóstico inicial.")
    notas_orador(s, "Contar el giro: el diagnóstico inicial estaba equivocado y "
                    "comprobarlo habilitó el dominio portuario. 20 segundos.")

    # 8 · Formulación -----------------------------------------------------------
    s = lamina(prs, "Formulación del problema", "", "Contexto")
    caja(s, 2.9, 5.9, 28.1, 5.4,
         [("¿Cómo desarrollar un producto de datos reproducible que integre información "
           "aduanera y portuaria de Buenaventura, permita monitorear y explicar "
           "indicadores agregados, pronostique únicamente aquellos con historia y calidad "
           "suficientes, comunique la incertidumbre de forma calibrada y no afirme "
           "relaciones que las fuentes no permiten sostener?", 18, True, BLANCO)],
         MARRON, MARRON, margen=1.1)
    _texto(s, 2.9, 11.9, 28.1, 2.0,
           [("La última condición de la pregunta admite una respuesta negativa, y el "
             "proyecto la mide en lugar de suponerla.", 14, False)],
           color=GRIS, align=PP_ALIGN.CENTER)
    notas_orador(s, "Leer la pregunta despacio. Subrayar que se admitía que la respuesta "
                    "fuera negativa. 20 segundos.")

    # 9 · Sistematización -------------------------------------------------------
    s = lamina(prs, "Sistematización del problema",
               "Diagrama de árbol: la pregunta central abierta en cuatro preguntas "
               "verificables", "Contexto")
    arbol(s, "¿Qué aporta integrar aduana y puerto en Buenaventura?", [
        ("¿Con qué llave?", "Qué tipo de unión admiten las fuentes: directa, agregada o "
                            "solo contextual."),
        ("¿Con qué calidad?", "Qué cobertura, continuidad y defectos tiene cada serie "
                              "una vez reconciliada."),
        ("¿Mejora el pronóstico?", "Si añadir variables portuarias reduce el error del "
                                   "pronóstico aduanero."),
        ("¿Qué no se puede responder?", "Qué preguntas quedan cerradas por ausencia de "
                                        "fuente pública."),
    ], y_raiz=5.9, y_ramas=9.8)
    notas_orador(s, "Cuatro preguntas, cada una con respuesta medida en el proyecto. 20 s.")

    # 10 · Objetivo general -----------------------------------------------------
    s = lamina(prs, "Objetivo general", "", "Contexto")
    caja(s, 2.9, 6.0, 28.1, 5.0,
         [("Desarrollar y validar un producto de datos reproducible que integre "
           "información aduanera y portuaria de Buenaventura para monitorear, explicar y "
           "pronosticar indicadores agregados, conservando la trazabilidad desde la "
           "fuente oficial hasta la salida de consumo.", 19, True, BLANCO)],
         MARRON, MARRON, margen=1.1)
    rejilla(s, 12.0, [[
        ("Trazable", "Cada cifra remite a un archivo de evidencia."),
        ("Validado", "Backtesting con ventana expansiva y líneas base."),
        ("Honesto", "Lo que la fuente no sostiene se declara no viable."),
    ]], alto=2.4, relleno=DURAZNO)
    notas_orador(s, "Subrayar la última línea del objetivo: solo lo que las fuentes "
                    "sostienen. 20 segundos.")

    # 11 · Objetivos específicos ------------------------------------------------
    s = lamina(prs, "Objetivos específicos",
               "Los ocho objetivos del documento, cada uno con el archivo que demuestra "
               "su cumplimiento", "Contexto")
    caja(s, 8.4, 5.2, 17.0, 1.3,
         [("Producto de datos integrado, trazable y validado", 13, True, BLANCO)],
         MARRON, MARRON)
    rejilla(s, 7.0, [
        [("1 · Evaluar la viabilidad de fuentes en cuatro dominios",
          "6 fuentes evaluadas: 2 integradas, 2 de contexto, 2 descartadas."),
         ("2 · Construir series mensuales reconciliadas",
          "173 meses aduaneros y 102 portuarios, continuos."),
         ("3 · Ejecutar las 52 preguntas y conservar la evidencia",
          "38 ejecutadas, 4 parciales y 10 no viables."),
         ("4 · Declarar el tipo de integración posible",
          "Integración agregada por mes; ninguna directa.")],
        [("5 · Determinar qué indicadores pueden pronosticarse",
          "4 elegibles de 7 evaluados."),
         ("6 · Medir si la integración mejora el pronóstico",
          f"No lo mejora: pasa de {c('WAPE historia propia')} % a "
          f"{c('WAPE historia propia mas puerto')} %."),
         ("7 · Comunicar la incertidumbre con cobertura medida",
          f"Cobertura empírica entre {c('Cobertura empirica minima')} % y "
          f"{c('Cobertura empirica maxima')} %."),
         ("8 · Documentar con evidencia lo no viable",
          f"{c('Preguntas no viables')} preguntas cerradas con búsqueda documentada.")],
    ], alto=3.7)
    notas_orador(s, "No leer los ocho: decir que cada uno tiene un archivo que lo prueba "
                    "y señalar el sexto, que es el que midió la hipótesis. 25 segundos.")

    # 12 · Justificación --------------------------------------------------------
    s = lamina(prs, "Justificación práctica",
               "Mapa mental: por qué el ejercicio tiene valor aunque parte de sus "
               "resultados sean negativos", "Contexto")
    mapa_conceptual(s, "¿Por qué\nhacerlo?", [
        ("Puerto crítico", "Buenaventura concentra buena parte del comercio del Pacífico "
                           "colombiano."),
        ("Fuentes abiertas", "Todo el insumo es público y la descarga es reproducible por "
                             "interfaz de programación."),
        ("Método replicable", "La arquitectura sirve para otros puertos con las mismas "
                              "fuentes."),
        ("Resultado negativo útil", "Medir que la integración no mejora el pronóstico "
                                    "evita construir sobre un supuesto falso."),
    ], cy=10.2)
    notas_orador(s, "Defender el resultado negativo como aporte, no como fracaso. 20 s.")

    # 13 · Alcance --------------------------------------------------------------
    s = lamina(prs, "Alcance",
               "Diagrama de Venn: lo que cada dominio aporta y lo que solo se ve al "
               "cruzarlos", "Contexto")
    venn(s,
         [("Aduana", 15, True), ("DANE IMPO", 10, False),
          (f"{c('Meses de la serie aduanera')} meses", 11, True),
          ("Valor CIF, peso neto,\norigen y capítulo", 9.5, False)],
         [("Puerto", 15, True), ("Supertransporte", 10, False),
          (f"{c('Meses de la serie portuaria')} meses", 11, True),
          ("Toneladas, tipo de carga\ny sociedad portuaria", 9.5, False)],
         [("Integración", 13, True), ("agregada por mes", 9.5, False),
          (f"{c('Meses integrados')} meses", 12, True),
          (c("Periodo integrado"), 9.5, False)],
         cy=10.9)
    nota(s, 15.4, "La intersección es temporal, no de registro: se cruzan meses, nunca una "
                  "declaración con un movimiento de carga.")
    notas_orador(s, "Dejar clarísimo que la unión es por mes. Es la limitación central. 20 s.")

    # 14 · Delimitación ---------------------------------------------------------
    s = lamina(prs, "Delimitación",
               "Tabla comparativa: qué entra, qué queda fuera y por qué", "Contexto")
    tabla(s, 1.9, 5.6, 30.1,
          ["Dimensión", "Dentro del proyecto", "Fuera del proyecto", "Razón"],
          [["Espacial", "Zona portuaria de Buenaventura y aduana 35",
            "Otros puertos del país", "El objeto de estudio es un solo nodo"],
           ["Temporal", f"Aduana {c('Periodo aduanero')}",
            "Antes de 2012 y después del último corte",
            "Cobertura efectiva de los microdatos"],
           ["Temática", "Importación, tráfico portuario y contexto",
            "Exportación como objeto de modelado", "La aduana 35 registra importación"],
           ["Operacional", "Toneladas y tipo de carga por sociedad",
            "Arribos, banderas, tiempos de atención",
            "Sin fuente pública tabular histórica"],
           ["Metodológica", "Modelos lineales regularizados con validación temporal",
            "Modelos de caja negra", "Muestra corta y exigencia de interpretabilidad"]],
          anchos=[3.2, 8.0, 7.4, 7.4], tam=10, alto_fila=1.32)
    nota(s, 14.6, f"De las 52 preguntas del catálogo, {c('Preguntas no viables')} quedaron "
                  "cerradas como no viables con la evidencia de dónde se buscó la fuente.")
    notas_orador(s, "La fila operacional es la que hay que defender: no es omisión, es "
                    "ausencia de fuente documentada. 20 segundos.")

    # 15 · Marco teórico y conceptual ------------------------------------------
    s = lamina(prs, "Marco teórico y conceptual",
               "Mapa conceptual: los conceptos que sostienen las decisiones técnicas",
               "Contexto")
    rejilla(s, 5.5, [
        [("Series de tiempo",
          "Autocorrelación, estacionalidad y tendencia como estructura a explotar.",
          f"ACF del CIF en el rezago 1: {c('ACF del CIF en rezago 1')}."),
         ("Validación temporal",
          "Ventana expansiva y pronóstico a un paso: nunca se entrena con el futuro.",
          "24 cortes de backtest sobre el valor CIF."),
         ("Líneas base",
          "Un modelo solo aporta si supera repetir el último valor.",
          f"Naive 1 sobre el CIF: {c('WAPE naive 1 sobre el CIF')} %.")],
        [("Incertidumbre",
          "Intervalos conformales calibrados sobre errores pasados, no supuestos.",
          f"Cobertura nominal {c('Cobertura nominal del intervalo')} %."),
         ("Concentración",
          "El índice de Herfindahl-Hirschman mide reparto, no capacidad instalada.",
          "Umbral de concentración: 2.500 puntos."),
         ("Trazabilidad",
          "Cada resultado nace de una pregunta del catálogo y deja un archivo.",
          "Catálogo cerrado de 52 preguntas.")],
    ], alto=4.0)
    notas_orador(s, "No definir cada concepto: decir para qué se usó cada uno. 25 segundos.")

    # 16 · Marco contextual y legal --------------------------------------------
    s = lamina(prs, "Marco contextual y legal",
               "Mapa de red institucional y jerarquía normativa de las fuentes",
               "Contexto")
    mapa_red(s, "Buenaventura", [
        ("DANE", "Microdatos de importación"),
        ("Supertransporte", "Tráfico portuario mensual"),
        ("DIMAR", "Boletines sin serie tabular"),
        ("Banco de la República", "Tasa de cambio representativa"),
        ("NOAA", "Índice oceánico El Niño"),
        ("Sociedades portuarias", "Operación de la zona"),
    ], cy=10.0, ry=3.5)
    _texto(s, 1.9, 14.35, 30.1, 0.6, [("Marco legal aplicado", 11, True)],
           color=MARRON)
    nota(s, 15.0, "1. Ley 1712 de 2014, transparencia y acceso a la información pública "
                  "  ·   2. Ley 1581 de 2012, protección de datos: el proyecto solo usa "
                  "agregados   ·   3. Licencia CC BY-SA 4.0 del conjunto portuario, "
                  "citada en cada uso   ·   4. NTC 1486 y APA séptima edición",
         color=TINTA)
    notas_orador(s, "Mencionar que no se tratan datos personales: todo es agregado. 20 s.")

    # ================== BLOQUE 2 · PRODUCTO DE DATOS (10 minutos) ==============

    # 17 · Metodología ----------------------------------------------------------
    s = lamina(prs, "Metodología",
               "Diagrama de flujo: el fundamento de ciencia de datos aplicado al proyecto",
               "Producto · 10 min")
    proceso(s, 5.5, [
        ("Comprensión\ndel negocio", "52 preguntas\ncerradas"),
        ("Comprensión\nde los datos", "Reconciliación\ny bitácora"),
        ("Preparación", "Cuatro capas\ncon hash"),
        ("Modelado", "Regresión\nregularizada"),
        ("Evaluación", "Walk-forward\ny líneas base"),
        ("Despliegue", "Cuadro de mando\nde solo lectura"),
    ], alto=3.1, tam=11)
    rejilla(s, 9.4, [[
        ("Sin fuga temporal",
         "Todas las variables entran rezagadas y los escaladores se ajustan dentro de "
         "cada corte, nunca sobre la serie completa."),
        ("Contra líneas base",
         "Cada modelo compite contra repetir el último valor, repetir el mismo mes del "
         "año anterior y la deriva."),
        ("Con ablación",
         "Se entrenan cinco configuraciones de variables para aislar qué grupo aporta y "
         "cuál estorba."),
    ]], alto=3.4)
    nota(s, 13.4, "La métrica principal es el error porcentual absoluto ponderado sobre "
                  "24 cortes de validación, no el ajuste sobre los datos de "
                  "entrenamiento.")
    notas_orador(s, "Este es el corazón metodológico. Detenerse en la ausencia de fuga "
                    "temporal. 45 segundos.")

    # 18 · El producto: propósito, público y valor ------------------------------
    s = lamina(prs, "El producto de datos",
               "Propósito, público objetivo y valor agregado", "Producto · 10 min")
    rejilla(s, 5.6, [[
        ("Propósito",
         "Leer en un solo lugar el comercio de importación y el tráfico portuario de "
         "Buenaventura, distinguiendo qué parte de un cambio viene del valor y qué parte "
         "del volumen físico."),
        ("Público objetivo",
         "Analistas de comercio exterior, equipos de planeación logística y "
         "observatorios académicos que hoy consultan las dos fuentes por separado."),
        ("Valor agregado",
         "Un contraste que ninguna fuente muestra sola, cifras con evidencia rastreable "
         "y un pronóstico que declara su margen de error en lugar de esconderlo."),
    ]], alto=5.0)
    rejilla(s, 11.3, [[
        ("Lo que sí hace",
         "Describe, compara dominios, pronostica lo que se puede pronosticar y publica "
         "el error de cada pronóstico."),
        ("Lo que no hace",
         "No estima congestión, tiempos de atención ni buques. No imputa valores "
         "ausentes: avisa de que faltan."),
    ]], alto=3.0, relleno=DURAZNO)
    notas_orador(s, "Aquí empieza la parte central. Nombrar al público concreto. 45 s.")

    # 19 · Ingeniería del producto ---------------------------------------------
    s = lamina(prs, "Ingeniería del producto",
               "Entorno de desarrollo, componentes y reutilización",
               "Producto · 10 min")
    tabla(s, 1.9, 5.5, 14.6, ["Componente", "Contenido"],
          [["Módulos de dominio", "Aduana, puerto e integración"],
           ["Módulos comunes", "15 heredados sin modificar de la versión anterior"],
           ["Suites de prueba", f"{c('Pruebas automatizadas')} pruebas automatizadas"],
           ["Revisión de estilo", "Sin hallazgos"],
           ["Cuadros de trabajo", "Entorno y análisis exploratorio"],
           ["Cuadro de mando", "Seis vistas de solo lectura"]],
          anchos=[6.0, 8.6], tam=10.5, alto_fila=0.95)
    rejilla(s, 5.5, [
        [("Entorno reproducible",
          "Dependencias fijadas en un archivo de requisitos; la ruta raíz se deriva del "
          "propio archivo, sin rutas absolutas en ningún módulo."),
         ],
        [("Ejecución local y en la nube",
          "El mismo paquete detecta si corre en el equipo o en un entorno de cuadernos en "
          "la nube y reapunta las rutas sin duplicar el sistema de carpetas."), ],
        [("Continuidad sobre lo que ya funcionaba",
          "La versión 5 no se reescribió: se conservaron los nombres y los módulos "
          "validados y solo se tocó lo que tenía un defecto demostrado."), ],
    ], alto=2.9, x0=17.2, ancho_total=14.8)
    nota(s, 14.9, "La reutilización es deliberada: métricas, intervalos, líneas base y "
                  "trazabilidad son agnósticos al dominio y sirven igual para toneladas "
                  "que para valor.")
    notas_orador(s, "Mostrar que hay ingeniería, no solo un cuaderno. 45 segundos.")

    # 20 · Arquitectura y trazabilidad -----------------------------------------
    s = lamina(prs, "Trazabilidad de punta a punta",
               "De la pregunta al archivo de evidencia y del archivo a la diapositiva",
               "Producto · 10 min")
    proceso(s, 5.6, [
        ("Pregunta", "Catálogo cerrado\nP01 a P52"),
        ("Ejecución", "Módulo o celda\ndel cuaderno"),
        ("Evidencia", "Archivo con\nhuella digital"),
        ("Cifra", "Catálogo de\ncifras verificadas"),
        ("Publicación", "Documento y\ncuadro de mando"),
    ], alto=2.9)
    rejilla(s, 9.4, [[
        ("Ejecutadas", c("Preguntas ejecutadas"), "con archivo de salida verificable"),
        ("Parciales", c("Preguntas parciales"), "ejecutadas con limitación declarada"),
        ("No viables", c("Preguntas no viables"), "con evidencia de dónde se buscó"),
        ("Sin justificar", "0", "ninguna quedó sin respuesta"),
    ]], alto=3.2, relleno=CREMA)
    nota(s, 13.4, f"El cuaderno de análisis produce {c('Figuras del cuaderno')} figuras y "
                  f"{c('Archivos de evidencia del cuaderno')} archivos de evidencia, y se "
                  "ejecuta de punta a punta sin errores. La construcción del documento "
                  "falla si alguien escribe una cifra que no esté en el catálogo.")
    notas_orador(s, "Este es el argumento de auditabilidad: ninguna cifra es escrita a "
                    "mano. 45 segundos.")

    # 21 · Hallazgo 1: la ablación ---------------------------------------------
    s = lamina(prs, "La integración no mejora el pronóstico",
               "Resultado central. Ablación multidominio sobre el valor CIF, 24 cortes "
               "de backtesting con ventana expansiva",
               "Producto · 10 min")
    tabla(s, 1.9, 5.6, 18.6,
          ["Conjunto de variables", "Variables", "Error (WAPE)", "Frente a A"],
          [[("B · historia propia + calendario", VERDE), "7",
            (f"{c('WAPE historia propia mas calendario')} %", VERDE), ("+0,207 pp", VERDE)],
           ["A · historia propia", "5", f"{c('WAPE historia propia')} %", "—"],
           ["D · A + contexto (tasa de cambio, El Niño)", "9",
            f"{c('WAPE contexto TRM y ONI')} %", "−0,051 pp"],
           ["E · integrado completo", "16", f"{c('WAPE integrado completo')} %",
            "−0,085 pp"],
           [("C · A + puerto", ROJO), "10",
            (f"{c('WAPE historia propia mas puerto')} %", ROJO), ("−0,493 pp", ROJO)],
           ["Base · repetir el último valor", "0", f"{c('WAPE naive 1 sobre el CIF')} %",
            "−0,874 pp"]],
          anchos=[9.0, 2.4, 3.6, 3.6], tam=10.5, alto_fila=0.98)
    rejilla(s, 5.6, [
        [("Por qué ocurre",
          "Ambas fuentes miden el mismo comercio subyacente y comparten el mismo rezago "
          "de publicación. El puerto no añade información que la historia del propio "
          "valor no contenga ya, y sí consume grados de libertad sobre una muestra "
          "corta."), ],
        [("Qué significa para el producto",
          "El valor de la integración no es predictivo: es explicativo. Sirve para "
          "responder si un mes cambió por valor o por volumen y para localizar por qué "
          "sociedad portuaria pasa la carga."), ],
    ], alto=3.4, x0=21.2, ancho_total=10.8, relleno=DURAZNO)
    nota(s, 13.2, "Las cinco configuraciones superan a las tres líneas base: el modelo "
                  "aporta sobre no hacer nada. Lo que no aporta es la integración. Se "
                  "reporta tal cual se midió.")
    notas_orador(s, "Momento clave de la sustentación. Decirlo sin rodeos: la hipótesis "
                    "de partida quedó refutada y se reporta. 60 segundos.")

    # 22 · Hallazgo 2: concentración -------------------------------------------
    s = lamina(prs, "Concentración por sociedad portuaria",
               "Lo que solo se ve integrando. Índice de Herfindahl-Hirschman en tres "
               "dimensiones del mismo comercio",
               "Producto · 10 min")
    rejilla(s, 5.6, [[
        ("Capítulo arancelario", mil("HHI capitulo arancelario"), "desconcentrado"),
        ("País de origen", mil("HHI pais de origen"), "desconcentrado"),
        ("Sociedad portuaria", mil("HHI sociedad portuaria"), "concentrado"),
    ]], alto=3.0, relleno=CREMA)
    _texto(s, 1.9, 9.1, 30.1, 1.6,
           ["Buenaventura importa mercancía variada desde orígenes variados, pero la "
            "mueve por muy pocas sociedades portuarias."],
           color=MARRON, align=PP_ALIGN.CENTER)
    tabla(s, 1.9, 10.2, 14.8, ["Sociedad portuaria", "Participación 2018–2026"],
          [[("Sociedad Portuaria Regional de Buenaventura", MARRON),
            (f"{c('Participacion de la mayor sociedad')} %", MARRON)],
           ["Sociedad Puerto Industrial Aguadulce", "25,28 %"],
           ["Sociedad Portuaria Terminal de Contenedores", "13,73 %"],
           ["Grupo Portuario", "6,01 %"],
           ["Compañía de Puertos Asociados", "3,62 %"]],
          anchos=[10.4, 4.4], tam=10, alto_fila=0.68)
    tabla(s, 17.4, 10.2, 14.6, ["Año", "Índice", "Sociedades que reportan"],
          [["2018", mil("HHI portuario 2018"), "5"],
           ["2023", mil("HHI portuario 2023"), "5"],
           ["2025", mil("HHI portuario 2025"), "5"],
           [("2026 · seis meses", GRIS), (mil("HHI portuario 2026 parcial"), GRIS),
            ("3", GRIS)],
           [("2025 recalculado con esas tres", GRIS), ("4.096", GRIS), ("3", GRIS)]],
          anchos=[6.6, 3.4, 4.6], tam=10, alto_fila=0.68)
    nota(s, 15.9, "El índice mide reparto de toneladas reportadas: no mide capacidad "
                  "instalada, utilización ni posibilidad real de sustitución. El salto "
                  "de 2026 se explica sobre todo por el cambio de cobertura del reporte.")
    notas_orador(s, "Advertir la limitación antes de que la pregunte el jurado. 60 s.")

    # 23 · Hallazgo 3: descomposición del tráfico -------------------------------
    s = lamina(prs, "La caída portuaria es transbordo",
               "Lo que esconde un total. Descomposición del tráfico de la zona portuaria "
               "por tipo de flujo",
               "Producto · 10 min")
    figura(s, "P23_carga_movilizada_en_la_zona_portuaria_de_buenaventura.png",
           1.9, 5.5, 18.4, 9.6)
    rejilla(s, 5.5, [
        [("Importación", f"{c('Variacion de toneladas de importacion')} %",
          "variación del periodo"), ],
        [("Exportación", f"{c('Variacion de toneladas de exportacion')} %",
          "variación del periodo"), ],
        [("Transbordo", f"{c('Variacion del transbordo')} %",
          "variación del periodo"), ],
        [("Total de la zona", f"{c('Variacion del total portuario')} %",
          "variación del periodo"), ],
    ], alto=2.15, x0=21.0, ancho_total=11.0, sep=0.35, relleno=CREMA)
    nota(s, 15.5, "El transbordo es carga que cambia de buque sin entrar ni salir del "
                  f"país: pesa {c('Transbordo sobre el total')} % del total movilizado y "
                  "sumarlo al comercio exterior infla la cifra. La caída del total no es "
                  "una caída del comercio.")
    notas_orador(s, "Ejemplo perfecto de por qué el total agregado engaña. 45 segundos.")

    # 24 · Hallazgo 4: valor y volumen -----------------------------------------
    s = lamina(prs, "Valor y volumen se separaron",
               "Descomposición logarítmica del crecimiento del valor CIF",
               "Producto · 10 min")
    figura(s, "P15_evolucion_aduanera_cif_peso_y_valor_unitario.png", 1.9, 5.5, 18.4, 9.6)
    rejilla(s, 5.5, [
        [("Crecimiento del valor CIF", f"{c('Crecimiento del CIF')} %",
          "en el periodo observado"), ],
        [("Aporte del volumen", f"{c('Aporte del volumen al crecimiento')} %",
          "de ese crecimiento"), ],
        [("Aporte del valor unitario", f"{c('Aporte del valor unitario al crecimiento')} %",
          "de ese crecimiento"), ],
    ], alto=2.6, x0=21.0, ancho_total=11.0, sep=0.4, relleno=CREMA)
    nota(s, 15.5, "Los dos aportes no suman exactamente cien: la diferencia es el residuo "
                  "de la descomposición y se declara. El valor CIF está en dólares "
                  "corrientes, de modo que parte de la variación es precio y no volumen, "
                  "y el cociente entre valor y kilogramo es un valor unitario declarado, "
                  "no un precio de mercado.")
    notas_orador(s, "Declarar el residuo y la naturaleza del valor unitario. 45 segundos.")

    # 25 · Qué se pronostica y qué no ------------------------------------------
    s = lamina(prs, "Qué se pronostica y qué no",
               "Un indicador que no supera la referencia trivial se presenta sin modelo",
               "Producto · 10 min")
    tabla(s, 1.9, 5.6, 30.1,
          ["Indicador", "Observaciones", "Mejor opción", "Error del modelo",
           "Repetir el último valor", "Decisión"],
          [["Valor CIF de importación", c("Meses de la serie aduanera"),
            "Historia propia + calendario",
            (f"{c('WAPE historia propia mas calendario')} %", VERDE),
            f"{c('WAPE naive 1 sobre el CIF')} %", ("se pronostica", VERDE)],
           ["Toneladas totales de la zona", c("Meses de la serie portuaria"),
            "Regresión regularizada",
            (f"{c('WAPE toneladas totales Ridge')} %", VERDE),
            f"{c('WAPE toneladas totales naive 1')} %", ("se pronostica", VERDE)],
           ["Carga contenerizada", c("Meses de la serie portuaria"),
            "Repetir el último valor",
            (f"{c('WAPE carga contenerizada Ridge')} %", ROJO),
            f"{c('WAPE carga contenerizada naive 1')} %",
            ("sin modelo", ROJO)],
           ["Contenedores equivalentes", "0", "—", "—", "—",
            ("la fuente no lo publica", GRIS)],
           ["Arribos y tiempos de atención", "0", "—", "—", "—",
            ("sin fuente pública", GRIS)]],
          anchos=[7.2, 3.0, 6.2, 4.4, 4.8, 4.5], tam=10, alto_fila=1.12)
    rejilla(s, 12.4, [[
        ("Intervalos calibrados, no supuestos",
         f"Cobertura nominal de {c('Cobertura nominal del intervalo')} % con cobertura "
         f"empírica observada entre {c('Cobertura empirica minima')} % y "
         f"{c('Cobertura empirica maxima')} % según el corte. La cobertura se mide, "
         "nunca se deriva del error."),
        ("Mejora frente a la referencia",
         f"{c('Mejora sobre naive 1')} % de mejora sobre repetir el último valor en el "
         f"valor CIF y {c('Mejora en toneladas totales')} % en toneladas totales. En "
         "carga contenerizada la mejora es negativa y por eso se retira el modelo."),
    ]], alto=3.2, relleno=DURAZNO)
    notas_orador(s, "El caso de la carga contenerizada es el que demuestra honestidad "
                    "metodológica. 60 segundos.")

    # 26 · Despliegue: el cuadro de mando --------------------------------------
    s = lamina(prs, "Despliegue del producto: el cuadro de mando",
               "Seis vistas de solo lectura construidas sobre los archivos de evidencia",
               "Producto · 10 min")
    rejilla(s, 5.6, [
        [("1 · Panorama", "Indicadores de cabecera de los dos dominios con su periodo y "
                          "su fuente."),
         ("2 · Aduana", "Serie del valor, del peso y del valor unitario declarado, con "
                        "filtros de periodo."),
         ("3 · Puerto", "Toneladas por tipo de carga y por sociedad portuaria.")],
        [("4 · Integración", "Comparación normalizada de los dos dominios y cobertura "
                             "mes a mes de la unión."),
         ("5 · Pronóstico", "Serie con intervalo, error del corte y comparación contra "
                            "las líneas base."),
         ("6 · Evidencia", "Estado de las 52 preguntas y descarga del archivo que "
                           "sostiene cada cifra.")],
    ], alto=3.4)
    rejilla(s, 12.9, [[
        ("Navegación", "Menú lateral por vista y filtros de periodo, dominio y sociedad "
                       "portuaria dentro de cada una."),
        ("Sin invención", "Si falta un archivo, la vista avisa de cuál falta y qué "
                          "comando lo genera. Nunca estima el valor ausente."),
        ("Solo lectura", "El cuadro de mando no escribe: consume lo que el proceso ya "
                         "produjo y verificó."),
    ]], alto=2.9, relleno=DURAZNO)
    notas_orador(s, "Enlazar con la demostración: nombrar las dos vistas que se van a "
                    "mostrar en vivo. 40 segundos.")

    # 27 · Demostración en vivo -------------------------------------------------
    s = lamina(prs, "Demostración en vivo", "Guion de dos minutos sobre el producto en "
                                            "ejecución", "Producto · 10 min")
    caja(s, 1.9, 5.5, 30.1, 1.5,
         [("streamlit run dashboard/app.py", 16, True, BLANCO)], TINTA, TINTA)
    proceso(s, 7.6, [
        ("1 · Panorama", "Abrir y leer\nlos indicadores"),
        ("2 · Descriptivo", "Filtrar el periodo\nen aduana y puerto"),
        ("3 · Integración", "Mostrar la unión\npor mes"),
        ("4 · Predictivo", "Serie con intervalo\ny error del corte"),
        ("5 · Evidencia", "Abrir el archivo\nque sostiene la cifra"),
    ], alto=2.9)
    rejilla(s, 11.5, [[
        ("Análisis descriptivo en vivo",
         "Se filtra un periodo en la vista de puerto y se muestra cómo cambia la "
         "composición por tipo de carga; se señala el transbordo separado del comercio "
         "exterior."),
        ("Análisis predictivo en vivo",
         "Se abre la vista de pronóstico, se muestra el intervalo alrededor de la "
         "proyección y se compara el error contra repetir el último valor en el mismo "
         "corte."),
    ]], alto=3.4, relleno=CREMA)
    nota(s, 15.4, "Plan de respaldo si falla la ejecución en vivo: capturas de las seis "
                  "vistas incorporadas al documento y los archivos de evidencia abiertos "
                  "directamente.")
    notas_orador(s, "Ensayar la demostración antes. Tener el plan de respaldo abierto en "
                    "otra pestaña. 120 segundos.")

    # ===================== BLOQUE 3 · CIERRE (2 minutos) =======================

    # 28 · Modelo de negocio ----------------------------------------------------
    s = lamina(prs, "Modelo de negocio",
               "Lienzo de modelo de negocio aplicado al producto de datos", "Cierre · 2 min")
    bloques = [
        (1.9, 5.5, 5.6, 6.2, "Aliados clave",
         "DANE · Supertransporte · Banco de la República · programa de Ciencia de Datos "
         "de la Universidad Libre"),
        (7.9, 5.5, 5.6, 3.0, "Actividades clave",
         "Ingesta reproducible, validación temporal y publicación de evidencia"),
        (7.9, 8.7, 5.6, 3.0, "Recursos clave",
         "Fuentes abiertas, código con pruebas y catálogo de preguntas"),
        (13.9, 5.5, 6.1, 6.2, "Propuesta de valor",
         "Leer aduana y puerto en un solo lugar, con cada cifra rastreable hasta su "
         "archivo y cada pronóstico acompañado de su error medido"),
        (20.4, 5.5, 5.6, 3.0, "Relación con el usuario",
         "Autoservicio guiado por el cuadro de mando y documentación abierta"),
        (20.4, 8.7, 5.6, 3.0, "Canales",
         "Cuadro de mando, repositorio público y documento académico"),
        (26.4, 5.5, 5.6, 6.2, "Segmentos de usuario",
         "Analistas de comercio exterior, planeación logística y observatorios "
         "académicos"),
        (1.9, 12.1, 14.5, 2.9, "Estructura de costos",
         "Cómputo local, almacenamiento de las fuentes crudas y tiempo de mantenimiento "
         "del proceso"),
        (17.0, 12.1, 15.0, 2.9, "Fuentes de ingreso",
         "Proyecto académico sin ánimo de lucro; el retorno es el ahorro de tiempo de "
         "consulta y la reutilización del método en otros puertos"),
    ]
    for x, y, an, al, titulo, cuerpo in bloques:
        relleno = DURAZNO if titulo == "Propuesta de valor" else BLANCO
        caja(s, x, y, an, al, [(titulo, 11, True, MARRON), (cuerpo, 9.5, False, TINTA)],
             relleno, MARRON, align=PP_ALIGN.LEFT, margen=0.26)
    notas_orador(s, "No leer los nueve bloques: nombrar propuesta de valor y segmentos. "
                    "25 segundos.")

    # 29 · Estrategias de comunicación -----------------------------------------
    s = lamina(prs, "Estrategias de comunicación",
               "Mapa de red: cómo llega el producto a cada público", "Cierre · 2 min")
    mapa_red(s, "Producto\nde datos", [
        ("Jurado y programa", "Documento y sustentación"),
        ("Comunidad académica", "Repositorio abierto y cuaderno reproducible"),
        ("Analistas del sector", "Cuadro de mando con filtros"),
        ("Semillero de investigación", "Método replicable en otros puertos"),
        ("Público general", "Resumen de hallazgos en lenguaje claro"),
    ], cy=10.1, ry=3.6)
    nota(s, 14.9, "Cada canal recibe el mismo hallazgo con distinto nivel de detalle, y "
                  "en todos los casos el hallazgo se puede rastrear hasta el archivo que "
                  "lo sostiene.")
    notas_orador(s, "20 segundos.")

    # 30 · Conclusiones ---------------------------------------------------------
    s = lamina(prs, "Conclusiones",
               "Lo que el proyecto puede afirmar con la evidencia recogida",
               "Cierre · 2 min")
    rejilla(s, 5.5, [
        [("1 · La integración es explicativa, no predictiva",
          "Añadir variables portuarias al pronóstico del valor CIF empeoró el error en "
          f"{c('WAPE historia propia mas puerto')} % frente a "
          f"{c('WAPE historia propia')} % de la historia propia. Se midió y se reporta."),
         ("2 · El contraste solo existe cruzando dominios",
          "Mercancía y orígenes variados frente a una concentración de "
          f"{mil('HHI sociedad portuaria')} puntos por sociedad portuaria: ninguna fuente "
          "por separado permite ver ese contraste.")],
        [("3 · Los totales agregados engañan",
          f"El total de la zona cae {c('Variacion del total portuario')} % mientras la "
          f"importación sube {c('Variacion de toneladas de importacion')} %: la caída es "
          "transbordo, que no es comercio exterior."),
         ("4 · Un producto honesto declara sus límites",
          f"{c('Preguntas no viables')} de las 52 preguntas quedaron cerradas con "
          "evidencia de por qué, y la carga contenerizada se presenta sin modelo porque "
          "ninguno superó la referencia trivial.")],
    ], alto=4.3)
    nota(s, 14.6, "El producto no puede afirmar nada sobre congestión, tiempos de "
                  "atención ni buques, y tampoco puede afirmar que dos sociedades "
                  "dejaran de operar: solo que no aparecen reportando en los seis meses "
                  "observados de 2026.")
    notas_orador(s, "Las cuatro conclusiones en cuarenta segundos. La cuarta es la que "
                    "define el tono del trabajo.")

    # 31 · Recomendaciones ------------------------------------------------------
    s = lamina(prs, "Recomendaciones",
               "Lista priorizada: qué haría más sólido el producto en su siguiente "
               "iteración", "Cierre · 2 min")
    prioridades = [
        ("Alta", "Validar las reglas de alerta con un usuario real del sector",
         "La propuesta de valor sigue siendo una hipótesis mientras nadie del público "
         "objetivo la use en una decisión.", ROJO),
        ("Alta", "Traducir los códigos de país y de capítulo arancelario a sus nombres "
                 "oficiales",
         "Hoy el cuadro de mando muestra códigos y eso limita la lectura de quien no "
         "conoce la nomenclatura.", ROJO),
        ("Media", "Incorporar una fuente institucional sobre la operación de las "
                  "sociedades portuarias",
         "Permitiría distinguir el cese de reporte del cese de operación, que hoy el "
         "proyecto no puede separar.", NARANJA),
        ("Media", "Medir la revisión de las cifras del DANE con una segunda descarga",
         "Con una sola descarga no se puede cuantificar cuánto cambian los meses ya "
         "publicados.", NARANJA),
        ("Baja", "Replicar la arquitectura en otro puerto con las mismas fuentes",
         "Comprobaría que el método es transferible y no está ajustado a Buenaventura.",
         GRIS),
    ]
    y = 5.5
    for prioridad, titulo, detalle, color in prioridades:
        etq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.9), Cm(y), Cm(2.6),
                                 Cm(1.75))
        _pinta(etq, color, color)
        _rellena_texto(etq, [(prioridad, 11, True, BLANCO)])
        caja(s, 4.7, y, 27.3, 1.75,
             [(titulo, 11.5, True, MARRON), (detalle, 9.5, False, TINTA)],
             BLANCO, MARRON, align=PP_ALIGN.LEFT, margen=0.3)
        y += 2.0
    notas_orador(s, "Nombrar solo las dos de prioridad alta. 25 segundos.")

    # 32 · Fuentes --------------------------------------------------------------
    s = lamina(prs, "Fuentes de información",
               "Todas las fuentes son públicas y su descarga es reproducible",
               "Cierre · 2 min")
    fuentes = [
        ("Datos", "DANE", "Microdatos de importación, aduana 35 de Buenaventura. "
                          f"{c('Registros de la aduana 35')} registros procesados, "
                          f"{c('Periodo aduanero')}."),
        ("Datos", "Superintendencia de Transporte",
         "Tráfico portuario mensual, conjunto 5r3g-zv5z de datos.gov.co, licencia "
         f"CC BY-SA 4.0. {c('Periodo portuario')}."),
        ("Contexto", "Banco de la República",
         "Tasa representativa del mercado, serie mensual."),
        ("Contexto", "NOAA", "Índice oceánico de El Niño, serie mensual."),
        ("Consultada", "DIMAR",
         "Boletines de tráfico marítimo. Sin serie tabular histórica descargable: por eso "
         "el dominio marítimo se declara no viable."),
        ("Norma", "NTC 1486 y APA séptima edición",
         "Presentación del documento y citación de fuentes."),
    ]
    y = 5.4
    for tipo, nombre, detalle in fuentes:
        icono = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.9), Cm(y), Cm(1.7), Cm(1.55))
        _pinta(icono, DORADO if tipo == "Datos" else CREMA, MARRON)
        _rellena_texto(icono, [(tipo[0], 13, True, MARRON)])
        _texto(s, 4.0, y + 0.1, 28.0, 1.5,
               [(nombre, 12, True, MARRON), (detalle, 10, False, TINTA)], interlinea=1.15)
        y += 1.62
    nota(s, 15.4, "Cada archivo descargado quedó registrado con su huella digital y su "
                  "fecha de consulta antes de ser procesado.")
    notas_orador(s, "No leerlas: señalar que todas son públicas y reproducibles. 20 s.")

    # 33 · Cierre ---------------------------------------------------------------
    s = _fondo(prs, "fondo_cierre.png")
    notas_orador(s, "Agradecer y abrir el espacio de preguntas.")

    salida = config.REPORTS / "Presentacion_Buenaventura_V5_INSTITUCIONAL.pptx"
    prs.save(salida)
    pd.DataFrame({"concepto": sorted(set(_USADAS))}).to_csv(
        S / "cifras_usadas_presentacion.csv", index=False)
    return salida


if __name__ == "__main__":
    ruta = construir()
    print(f"presentación guardada en {ruta}")
    print(f"cifras verificadas utilizadas: {len(set(_USADAS))}")
